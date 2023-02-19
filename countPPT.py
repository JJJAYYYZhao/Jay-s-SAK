import os
import os.path
import pptx

total = 0

import pdfplumber


def pptCount(path):
    # 声明全局变量
    global total
    # 遍历文件夹中所有文件或子文件夹
    for subPath in os.listdir(path):
        subPath = os.path.join(path, subPath)
        if os.path.isdir(subPath):
            # 递归遍历子文件夹
            pptCount(subPath)
        elif subPath.endswith('.pptx'):
            # 显示正在处理的文件名
            print(subPath.split("\\")[-1])
            # 统计幻灯片数量
            presentation = pptx.Presentation(subPath)
            print(len(presentation.slides))
            total += len(presentation.slides)
        elif subPath.endswith('.pdf'):
            # 显示正在处理的文件名
            print(subPath.split("\\")[-1])
            f = pdfplumber.open(subPath)
            print(len(f.pages))
            total += len(f.pages)


# xxx中填入你要统计的多个ppt文件的父文件夹所在位置，例如'F:\\教学课件\\Python程序设计（第二版）'
# pptCount('C:\\Users\\16551\Documents\\WeChat Files\\wxid_su61s6y9dw1p22\\FileStorage\\File\\2023-02\\ppt\\ppt')
pptCount('C:\\Users\\16551\\Desktop\\大数据\\ppt')
print('总页数: ', total)
